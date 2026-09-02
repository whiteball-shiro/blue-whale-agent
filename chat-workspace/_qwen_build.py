# -*- coding: utf-8 -*-
"""本地 qwen 文档构建工具：由 local-qwen.mjs 调用。
用法: python _qwen_build.py <input.json>
input.json: {"kind": "docx|pdf|pptx", "out": "绝对路径", "text": "内容", "slides": 1}
"""
import json
import sys
from pathlib import Path

payload = json.loads(Path(sys.argv[1]).read_text(encoding="utf-8"))
kind = payload["kind"]
out = payload["out"]
text = payload.get("text", "")
slides = max(1, int(payload.get("slides") or 1))

if kind == "docx":
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
    from docx.shared import Inches, Pt

    doc = Document()
    sec = doc.sections[0]
    sec.page_width = Inches(8.5)
    sec.page_height = Inches(11)
    sec.top_margin = Inches(1)
    sec.right_margin = Inches(1)
    sec.bottom_margin = Inches(1)
    sec.left_margin = Inches(1)
    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)
    pf = normal.paragraph_format
    pf.alignment = WD_ALIGN_PARAGRAPH.LEFT
    pf.space_before = Pt(0)
    pf.space_after = Pt(6)
    pf.line_spacing = 1.25
    pf.line_spacing_rule = WD_LINE_SPACING.MULTIPLE
    doc.add_paragraph(text)
    doc.save(out)

elif kind == "pdf":
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import Paragraph, SimpleDocTemplate

    pdfmetrics.registerFont(TTFont("Calibri", r"C:\Windows\Fonts\calibri.ttf"))
    style = ParagraphStyle(
        "Body", fontName="Calibri", fontSize=11, leading=11 * 1.25,
        spaceBefore=0, spaceAfter=6, alignment=0,
    )
    doc = SimpleDocTemplate(
        out, pagesize=letter,
        leftMargin=1 * inch, rightMargin=1 * inch,
        topMargin=1 * inch, bottomMargin=1 * inch,
        title=text,
    )
    doc.build([Paragraph(text, style)])

elif kind == "pptx":
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt

    slides_text = payload.get("slides_text")
    per_slide = isinstance(slides_text, list) and len(slides_text) > 0

    if per_slide:
        # 新版式：每页内容不同（第一行标题，- 开头行为要点），顶部标题栏 + 要点 + 页码
        pages = [str(x) for x in slides_text]
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        NAVY = RGBColor(0x1F, 0x38, 0x64)
        DARK = RGBColor(0x33, 0x33, 0x33)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        GRAY = RGBColor(0x8A, 0x8A, 0x8A)
        FONT = "Microsoft YaHei"

        def set_cn_font(run, name):
            run.font.name = name
            rPr = run._r.get_or_add_rPr()
            for tag in ("a:ea", "a:cs"):
                el = rPr.find(qn(tag))
                if el is None:
                    el = rPr.makeelement(qn(tag), {})
                    rPr.append(el)
                el.set("typeface", name)

        for idx, page in enumerate(pages, start=1):
            lines = [ln.strip() for ln in str(page).splitlines() if ln.strip()]
            title = lines[0] if lines else "幻灯片"
            bullets = []
            for ln in lines[1:]:
                if ln.startswith("- "):
                    bullets.append(ln[2:].strip())
                elif ln.startswith("-"):
                    bullets.append(ln[1:].strip())
                else:
                    bullets.append(ln)

            slide = prs.slides.add_slide(blank)
            band = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.35)
            )
            band.fill.solid()
            band.fill.fore_color.rgb = NAVY
            band.line.fill.background()
            band.shadow.inherit = False

            tb = slide.shapes.add_textbox(
                Inches(0.6), Inches(0.16), prs.slide_width - Inches(1.2), Inches(1.05)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = title
            r.font.size = Pt(34)
            r.font.bold = True
            r.font.color.rgb = WHITE
            set_cn_font(r, FONT)

            body = slide.shapes.add_textbox(
                Inches(0.9), Inches(1.75), Inches(11.5), Inches(5.3)
            )
            bf = body.text_frame
            bf.word_wrap = True
            for j, item in enumerate(bullets):
                p = bf.paragraphs[0] if j == 0 else bf.add_paragraph()
                p.space_after = Pt(16)
                r = p.add_run()
                r.text = "•  " + item
                r.font.size = Pt(22)
                r.font.color.rgb = DARK
                set_cn_font(r, FONT)

            pn = slide.shapes.add_textbox(
                prs.slide_width - Inches(1.3),
                prs.slide_height - Inches(0.55),
                Inches(1.0),
                Inches(0.4),
            )
            pfp = pn.text_frame.paragraphs[0]
            pfp.alignment = PP_ALIGN.RIGHT
            pr_ = pfp.add_run()
            pr_.text = f"{idx} / {len(pages)}"
            pr_.font.size = Pt(12)
            pr_.font.color.rgb = GRAY
            set_cn_font(pr_, FONT)

        prs.save(out)
    else:
        # 原有逻辑：每页显示相同文本（向后兼容）
        prs = Presentation()
        blank = prs.slide_layouts[6]
        for _ in range(slides):
            slide = prs.slides.add_slide(blank)
            tb = slide.shapes.add_textbox(Inches(1), Inches(2.75), Inches(8), Inches(2))
            tf = tb.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            run.font.size = Pt(72)
            run.font.bold = True
            run.font.name = "Calibri"
        prs.save(out)

elif kind == "xlsx":
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"

    rows = payload.get("rows")
    if not isinstance(rows, list):
        # 兼容纯文本：按行拆分，遇 Tab 或逗号则分列
        rows = []
        for ln in str(text).splitlines():
            ln = ln.strip()
            if not ln:
                continue
            if "\t" in ln:
                rows.append([c.strip() for c in ln.split("\t")])
            elif "," in ln:
                rows.append([c.strip() for c in ln.split(",")])
            else:
                rows.append([ln])

    has_header = bool(payload.get("header", True))
    fill = PatternFill(start_color="FF1F3864", end_color="FF1F3864", fill_type="solid")
    head_font = Font(bold=True, color="FFFFFFFF")
    thin = Side(style="thin", color="FFCCCCCC")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    for r_idx, row in enumerate(rows):
        cells = []
        if isinstance(row, (list, tuple)):
            cells = list(row)
        elif isinstance(row, dict):
            cells = [row.get("v", row.get("value", ""))]
        else:
            cells = [row]
        for c_idx, val in enumerate(cells):
            cell = ws.cell(row=r_idx + 1, column=c_idx + 1, value=val)
            cell.border = border
            cell.alignment = Alignment(vertical="center")
            if has_header and r_idx == 0:
                cell.font = head_font
                cell.fill = fill
                cell.alignment = Alignment(vertical="center", horizontal="center")

    # 列宽自适应（粗略；中文按 2 计）
    for col in ws.columns:
        max_len = 0
        for cell in col:
            v = cell.value
            if v is not None:
                s = str(v)
                w = sum(2 if ord(ch) > 127 else 1 for ch in s)
                max_len = max(max_len, w)
        if max_len:
            ws.column_dimensions[col[0].column_letter].width = min(max(max_len * 1.2, 8), 50)

    wb.save(out)

else:
    raise SystemExit("unknown kind: " + kind)

print("OK")
